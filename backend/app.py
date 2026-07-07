from flask import Flask, jsonify, request, session
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
from dotenv import load_dotenv
from models import db, User
import os
from routes.auth import auth_bp
from routes.sessions import sessions_bp
from routes.user import user_bp
from routes.bookmarks import bookmarks_bp
from routes.courses import courses_bp
from datetime import date, timedelta


load_dotenv()

app = Flask(__name__)
CORS(app, supports_credentials=True, origins=['http://localhost:5173', 'https://adapt-learn-eta.vercel.app'])


app.config['SECRET_KEY'] = os.getenv('SECRET_KEY')

DATABASE_URL = os.getenv('DATABASE_URL', 'sqlite:///adaptlearn.db')
if DATABASE_URL.startswith('postgres://'):
    DATABASE_URL = DATABASE_URL.replace('postgres://', 'postgresql://', 1)
app.config['SQLALCHEMY_DATABASE_URI'] = DATABASE_URL
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db.init_app(app)

@app.route('/api/hello')
def hello():
    return jsonify({'message': 'AdaptLearn is alive'})

@app.route('/api/explain', methods=['POST'])
def explain():
    data = request.json
    if not data or not data.get('topic'):
        return jsonify({'error': 'Topic is required'}), 400
    if not data.get('style'):
        return jsonify({'error': 'Style is required'}), 400

    topic = data['topic']
    style = data['style']
    education_level = data.get('education_level')
    user_id = session.get('user_id')

    try:
        from services.claude_service import get_explanation
        explanation_text = get_explanation(topic, style, education_level)

        # save to database if logged in
        session_id = None
        explanation_id = None

        if user_id:
            from models import LearningSession, Explanation
            from datetime import datetime

            # create session
            learning_session = LearningSession(
                user_id=user_id,
                topic=topic,
                initial_style=style,
                started_at=datetime.utcnow()
            )
            db.session.add(learning_session)
            db.session.flush()

            # create explanation
            explanation_record = Explanation(
                session_id=learning_session.id,
                style=style,
                prompt_sent=f"{education_level} | {style} | {topic}",
                response_text=explanation_text,
                attempt_number=1,
                model_used='claude-haiku-4-5-20251001'
            )
            db.session.add(explanation_record)
            db.session.commit()

            session_id = learning_session.id
            explanation_id = explanation_record.id

        return jsonify({
            'explanation': explanation_text,
            'style': style,
            'topic': topic,
            'session_id': session_id,
            'explanation_id': explanation_id,
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
@app.route('/api/explanations/<int:explanation_id>/rag', methods=['PATCH'])
def rate_explanation(explanation_id):
    data = request.json
    rating = data.get('rating')

    if rating not in ['red', 'amber', 'green']:
        return jsonify({'error': 'Invalid rating'}), 400

    from models import Explanation
    explanation = Explanation.query.get(explanation_id)
    if not explanation:
        return jsonify({'error': 'Explanation not found'}), 404

    explanation.rag_rating = rating
    db.session.commit()

    return jsonify({
        'success': True,
        'next_action': 'quiz' if rating == 'green' else 'retry'
    })

@app.route('/api/quiz', methods=['POST'])
def quiz():
    data = request.json
    if not data or not data.get('explanation_text'):
        return jsonify({'error': 'Explanation text is required'}), 400
    explanation_text = data['explanation_text']
    num_questions = data.get('num_questions', 3)
    education_level = data.get('education_level')
    try:
        from services.claude_service import generate_quiz
        questions = generate_quiz(explanation_text, num_questions, education_level)
        return jsonify({'questions': questions, 'num_questions': num_questions})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/quiz/submit', methods=['POST'])
def submit_quiz():
    data = request.json
    explanation_id = data.get('explanation_id')
    answers = data.get('answers', [])
    questions = data.get('questions', [])

    if not explanation_id:
        return jsonify({'error': 'explanation_id required'}), 400

    try:
        from models import QuizResult, Explanation, LearningSession, UserStats
        import json

        correct_count = 0

        for i, question in enumerate(questions):
            user_answer = answers[i] if i < len(answers) else None
            is_correct = user_answer == question.get('correct_index')
            if is_correct:
                correct_count += 1

            result = QuizResult(
                explanation_id=explanation_id,
                question=question['question'],
                options=json.dumps(question['options']),
                correct_index=question['correct_index'],
                user_answer_index=user_answer,
                correct=is_correct
            )
            db.session.add(result)

        score_pct = round((correct_count / len(questions)) * 100) if questions else 0

        # update session score
        explanation = Explanation.query.get(explanation_id)
        if explanation:
            learning_session = LearningSession.query.get(explanation.session_id)
            if learning_session:
                learning_session.final_quiz_score = score_pct
                learning_session.completed_at = __import__('datetime').datetime.utcnow()

                # award XP
                xp_earned = max(5, 10 + max(0, score_pct - 60))
                learning_session.xp_earned = xp_earned

                # update user stats
                user_stats = UserStats.query.filter_by(
                    user_id=learning_session.user_id
                ).first()
                if user_stats:
                    user_stats.xp += xp_earned
                    user_stats.total_sessions += 1
                    user_stats.total_topics += 1
                
                # after updating user_stats xp


                if user_stats:
                    user_stats.xp += xp_earned
                    user_stats.total_sessions += 1
                    user_stats.total_topics += 1

    # update streak
                today = date.today()
                last = user_stats.last_active.date() if user_stats.last_active else None
                if last is None or last < today - timedelta(days=1):
                    user_stats.streak = 1
                elif last == today - timedelta(days=1):
                    user_stats.streak += 1
                # if last == today, streak stays same — already counted
                user_stats.last_active = __import__('datetime').datetime.utcnow()

        db.session.commit()

        return jsonify({
            'score_pct': score_pct,
            'correct_count': correct_count,
            'total': len(questions),
            'xp_earned': xp_earned if questions else 5
        })


    except Exception as e:
        return jsonify({'error': str(e)}), 500
@app.route('/api/extract', methods=['POST'])
def extract():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({'error': 'Not logged in'}), 401

    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    filename = file.filename
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    if ext not in ['pdf', 'pptx', 'txt', 'md']:
        return jsonify({'error': 'Unsupported file type'}), 400

    import tempfile
    import os

    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{ext}') as tmp:
        file.save(tmp.name)
        tmp_path = tmp.name

    try:
        # extract text
        text = ''
        if ext == 'pdf':
            import fitz
            doc = fitz.open(tmp_path)
            text = '\n'.join(page.get_text() for page in doc)
            doc.close()  # ← add this line to close before deleting
            if len(text.strip()) < 100:
                return jsonify({'error': 'PDF appears to be scanned — no text layer found'}), 400
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(tmp_path)
            text = '\n'.join(
                tf.text for sl in prs.slides
                for sh in sl.shapes if sh.has_text_frame
                for tf in sh.text_frame.paragraphs
            )
        else:
            with open(tmp_path, encoding='utf-8') as f:
                text = f.read()

        # extract topics with Claude
        import anthropic
        client = anthropic.Anthropic(api_key=os.getenv('ANTHROPIC_API_KEY'))
        prompt = f"""From this educational document, identify every distinct concept or topic a student would need to understand.
Return ONLY a JSON array of short topic title strings, ordered logically.
No markdown, no preamble. Example: ["Multi-store model","Working memory","Encoding"]
Document: {text[:4000]}"""

        message = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=1000,
            messages=[{'role': 'user', 'content': prompt}]
        )

        import re
        import json
        raw = message.content[0].text
        clean = re.sub(r'```json|```', '', raw).strip()
        start = clean.find('[')
        end = clean.rfind(']') + 1
        topics = json.loads(clean[start:end]) if start != -1 else []

        return jsonify({
            'topics': topics,
            'doc_text': text[:8000],
            'filename': filename
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
    
app.register_blueprint(auth_bp)
app.register_blueprint(sessions_bp)
app.register_blueprint(user_bp)
app.register_blueprint(bookmarks_bp)
app.register_blueprint(courses_bp)

with app.app_context():
    db.create_all()
    print("Tables created!")

if __name__ == '__main__':
    app.run(debug=True)